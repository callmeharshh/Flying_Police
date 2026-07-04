"""Generate architecture PPT for Flying Police."""
import os
import sys

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0D, 0x11, 0x17)   # near-black
CARD_DARK   = RGBColor(0x16, 0x1B, 0x22)   # dark card
CARD_MID    = RGBColor(0x1C, 0x27, 0x33)   # mid card
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)   # sky-400
ACCENT_RED  = RGBColor(0xF8, 0x71, 0x71)   # red-400
ACCENT_GRN  = RGBColor(0x34, 0xD3, 0x99)   # emerald-400
ACCENT_YLW  = RGBColor(0xFB, 0xBF, 0x24)   # amber-400
ACCENT_PRP  = RGBColor(0xA7, 0x8B, 0xFA)   # violet-400
ACCENT_ORG  = RGBColor(0xFB, 0x92, 0x3C)   # orange-400
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY        = RGBColor(0x94, 0xA3, 0xB8)   # slate-400

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def add_rounded_rect(slide, x, y, w, h, fill: RGBColor, line_color: RGBColor = None, line_width=Pt(1)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    # Adjust corner radius
    adj = shape.adjustments
    if adj:
        adj[0] = 0.05
    return shape


def add_text_box(slide, text, x, y, w, h, font_size=Pt(11), bold=False,
                 color=WHITE, align=PP_ALIGN.CENTER, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_arrow(slide, x1, y1, x2, y2, color: RGBColor = GREY, width=Pt(1.5)):
    """Add a connector arrow between two points (in inches)."""
    from pptx.util import Inches
    connector = slide.shapes.add_connector(
        pptx.enum.shapes.PP_CONNECTOR_TYPE.STRAIGHT if False else 1,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_component(slide, x, y, w, h, title, subtitle, accent: RGBColor, icon=""):
    """Draw a component card with accent top bar."""
    # Card background
    card = add_rounded_rect(slide, x, y, w, h, CARD_MID, line_color=accent, line_width=Pt(1.2))
    # Top accent bar
    bar = add_rect(slide, x, y, w, 0.06, accent)
    # Icon + title
    label = f"{icon}  {title}" if icon else title
    add_text_box(slide, label, x, y + 0.08, w, 0.28,
                 font_size=Pt(10.5), bold=True, color=accent)
    # Subtitle
    add_text_box(slide, subtitle, x + 0.05, y + 0.35, w - 0.1, h - 0.42,
                 font_size=Pt(8.5), color=GREY, align=PP_ALIGN.LEFT, wrap=True)
    return card


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    # Accent bar
    add_rect(slide, 0, 3.2, 13.33, 0.06, ACCENT_BLUE)

    add_text_box(slide, "FLYING POLICE",
                 0.5, 1.8, 12.33, 1.0,
                 font_size=Pt(36), bold=True, color=WHITE)
    add_text_box(slide, "System Architecture & Component Overview",
                 0.5, 2.8, 12.33, 0.6,
                 font_size=Pt(18), color=ACCENT_BLUE)
    add_text_box(slide,
                 "Real-time video analysis  ·  VLM captioning  ·  Rule-based alerting  ·  LangChain agent  ·  Dual-DB indexing  ·  Telegram notifications",
                 0.5, 3.5, 12.33, 0.5,
                 font_size=Pt(11), color=GREY)

    # Tech stack pills
    pills = [
        ("OpenCV", ACCENT_BLUE), ("BLIP VLM", ACCENT_GRN),
        ("LangChain", ACCENT_PRP), ("GPT-4o-mini", ACCENT_YLW),
        ("ChromaDB", ACCENT_ORG), ("SQLite", ACCENT_RED),
        ("Telegram Bot", ACCENT_BLUE),
    ]
    px = 1.0
    for label, color in pills:
        w = len(label) * 0.11 + 0.4
        pill = add_rounded_rect(slide, px, 4.4, w, 0.35, CARD_MID, line_color=color)
        add_text_box(slide, label, px, 4.42, w, 0.32, font_size=Pt(9.5), color=color)
        px += w + 0.18

    add_text_box(slide, "github.com/callmeharshh/Flying-Police", 0.5, 6.8, 12.33, 0.4,
                 font_size=Pt(10), color=GREY)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2 — Full Pipeline Flow (clean numbered lane)
# ─────────────────────────────────────────────────────────────────────────────
def add_step(slide, num, title, body, accent, x, y, w=1.7, h=1.35):
    """Single numbered pipeline step card."""
    # Card
    add_rounded_rect(slide, x, y, w, h, CARD_MID, line_color=accent, line_width=Pt(1.5))
    # Number bubble (top-left corner)
    bubble = add_rounded_rect(slide, x + 0.08, y + 0.08, 0.38, 0.38, accent)
    add_text_box(slide, str(num), x + 0.08, y + 0.07, 0.38, 0.38,
                 font_size=Pt(11), bold=True, color=BG)
    # Title
    add_text_box(slide, title, x + 0.52, y + 0.08, w - 0.6, 0.38,
                 font_size=Pt(9.5), bold=True, color=accent, align=PP_ALIGN.LEFT)
    # Body
    add_text_box(slide, body, x + 0.12, y + 0.5, w - 0.2, h - 0.56,
                 font_size=Pt(7.8), color=GREY, align=PP_ALIGN.LEFT, wrap=True)


def add_horiz_arrow(slide, x, y, length=0.22, color=GREY):
    add_arrow(slide, x, y, x + length, y, color=color, width=Pt(2))


def add_vert_arrow(slide, x, y, length=0.28, color=GREY):
    add_arrow(slide, x, y, x, y + length, color=color, width=Pt(2))


def slide_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    # ── Title ──
    add_text_box(slide, "END-TO-END PIPELINE", 0.3, 0.1, 12.73, 0.38,
                 font_size=Pt(15), bold=True, color=WHITE)
    add_text_box(slide, "How a raw video frame becomes a security alert",
                 0.3, 0.46, 12.73, 0.28, font_size=Pt(9.5), color=GREY)

    # ── Section label: PROCESSING PIPELINE ──
    add_rect(slide, 0.28, 0.82, 0.06, 1.55, ACCENT_BLUE)
    add_text_box(slide, "PROCESSING", 0.36, 0.82, 0.9, 0.95,
                 font_size=Pt(7.5), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.LEFT)

    # ── Row 1: steps 1–7 across the slide ──
    sw, sh = 1.7, 1.35          # step card width / height
    gap    = 0.18               # gap between cards
    start_x = 1.3
    row_y   = 0.82

    steps = [
        (1, "Video Input",        "MP4 file\nAny resolution\nAny FPS",                                  ACCENT_BLUE),
        (2, "Frame Sampler",      "Every Nth frame\nSkips redundant\nframes (speed)",                   ACCENT_BLUE),
        (3, "MOG2 Warmup",        "First 5 frames\nbuild background\nmodel (calibrate)",                ACCENT_GRN),
        (4, "Foreground Detect",  "Moving blobs found\nFiltered by size\n(2000px²–40%)",                ACCENT_GRN),
        (5, "BLIP Caption",       "Crop sent to VLM\nReturns plain-text\ndescription",                  ACCENT_PRP),
        (6, "Activity Infer",     "Caption keywords\n→ entering /\nloitering / stationary",             ACCENT_YLW),
        (7, "Alert Rules",        "RULE-01 → 05\nDeterministic\nfires before LLM",                     ACCENT_RED),
    ]

    for i, (num, title, body, color) in enumerate(steps):
        x = start_x + i * (sw + gap)
        add_step(slide, num, title, body, color, x, row_y, sw, sh)
        if i < len(steps) - 1:
            ax = x + sw
            ay = row_y + sh / 2
            add_horiz_arrow(slide, ax, ay, length=gap, color=GREY)

    # ── Step 8: LangChain (below step 7, turns down) ──
    s7_x = start_x + 6 * (sw + gap)
    s8_x = s7_x
    s8_y = row_y + sh + 0.42

    # Down arrow from step 7
    add_vert_arrow(slide, s7_x + sw / 2, row_y + sh, length=0.38, color=ACCENT_ORG)

    add_step(slide, 8, "LangChain Agent",
             "ReAct + GPT-4o-mini\nContext-aware log\nk=10 memory window",
             ACCENT_ORG, s8_x, s8_y, sw, sh)

    # ── Section label: OUTPUTS ──
    add_rect(slide, 0.28, s8_y, 0.06, 1.55, ACCENT_GRN)
    add_text_box(slide, "OUTPUTS", 0.36, s8_y + 0.3, 0.9, 0.6,
                 font_size=Pt(7.5), bold=True, color=ACCENT_GRN, align=PP_ALIGN.LEFT)

    # ── Output row: 4 boxes aligned under steps 1–6 ──
    out_y = s8_y
    outputs = [
        ("SQLite\nEvent Store",    "Events + alerts logged\nTime & severity queries",         ACCENT_RED,  start_x),
        ("ChromaDB\nFrame Index",  "Vector embedding per frame\nSemantic NL search",          ACCENT_ORG,  start_x + 1 * (sw + gap) * 1.6),
        ("Query Engine",           "Routes questions to\nSQLite or ChromaDB",                 ACCENT_BLUE, start_x + 2 * (sw + gap) * 1.6),
        ("Telegram\nNotification", "HIGH alerts only\nAnnotated frame photo",                 ACCENT_GRN,  start_x + 3 * (sw + gap) * 1.6),
    ]

    for title, body, color, ox in outputs:
        add_component(slide, ox, out_y, sw + 0.2, sh, title, body, color)

    # ── Arrows: step 7 → SQLite, step 7 → Telegram ──
    # Alert Rules → SQLite (long diagonal goes left)
    add_arrow(slide, s7_x + sw / 2, row_y + sh,
              start_x + (sw + 0.2) / 2, out_y,
              color=ACCENT_RED, width=Pt(1.2))

    # Alert Rules → Telegram
    add_arrow(slide, s7_x + sw / 2, s8_y,
              start_x + 3 * (sw + gap) * 1.6 + (sw + 0.2) / 2, out_y,
              color=ACCENT_GRN, width=Pt(1.2))

    # BLIP → ChromaDB
    blip_x = start_x + 4 * (sw + gap) + sw / 2
    add_vert_arrow(slide, blip_x, row_y + sh,
                   length=(out_y - row_y - sh + 0.28), color=ACCENT_PRP)
    add_arrow(slide, blip_x, out_y,
              start_x + 1 * (sw + gap) * 1.6 + (sw + 0.2) / 2, out_y,
              color=ACCENT_PRP, width=Pt(1.2))

    # LangChain → SQLite (short arrow left)
    add_arrow(slide, s8_x, s8_y + sh / 2,
              start_x + (sw + 0.2), out_y + sh / 2,
              color=ACCENT_ORG, width=Pt(1.2))

    # ── Legend ──
    legend = [("Video → Detection → Caption → Rules → Agent", GREY),
              ("→ SQLite (events/alerts)  +  ChromaDB (semantic search)  +  Telegram (HIGH alerts)", GREY)]
    for i, (txt, color) in enumerate(legend):
        add_text_box(slide, txt, 0.28, 6.75 + i * 0.28, 12.77, 0.28,
                     font_size=Pt(7.8), color=color, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 3 — VLM Layer Detail
# ─────────────────────────────────────────────────────────────────────────────
def slide_vlm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_text_box(slide, "VLM LAYER — Object Detection & Captioning", 0.3, 0.15, 12.73, 0.4,
                 font_size=Pt(14), bold=True, color=WHITE)

    steps = [
        ("1. Frame Sampled", "Every Nth raw frame\nexits the video stream", ACCENT_BLUE),
        ("2. MOG2 Applied", "Background model built\nover warmup frames.\nForeground mask extracted", ACCENT_GRN),
        ("3. Contour Filter", "Min area: 2000 px²\nMax area: 40% of frame\nShadow pixels removed", ACCENT_YLW),
        ("4. Object Crop", "Bounding box + 10px\npadding cropped\nConverted to PIL Image", ACCENT_ORG),
        ("5. BLIP Caption", "Salesforce/blip-image-\ncaptioning-base\nRuns on CPU / MPS", ACCENT_PRP),
        ("6. Keyword Extract", "Objects: person / vehicle\nColors: red / blue / …\nBrand names filtered", ACCENT_RED),
    ]

    for i, (title, body, color) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.25
        y = 1.1 + row * 2.5
        add_component(slide, x, y, 3.9, 2.1, title, body, color)
        if col < 2:
            add_arrow(slide, x + 3.9, y + 1.05, x + 4.25, y + 1.05, color=GREY)

    add_text_box(slide,
                 "Output: VLMResult(raw_description, objects[ ], colors[ ], bbox, bbox_relative)",
                 0.4, 6.0, 12.5, 0.4, font_size=Pt(10), color=ACCENT_BLUE,
                 align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4 — Alert Rules
# ─────────────────────────────────────────────────────────────────────────────
def slide_rules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_text_box(slide, "DETERMINISTIC ALERT RULES ENGINE", 0.3, 0.15, 12.73, 0.4,
                 font_size=Pt(14), bold=True, color=WHITE)
    add_text_box(slide, "Fires before LLM — zero hallucination risk on alert logic", 0.3, 0.52, 12.73, 0.3,
                 font_size=Pt(10), color=GREY)

    rules = [
        ("RULE-01", "After-Hours Person",
         "Person detected between 22:00–06:00\nFires once per tracked person\nSeverity: HIGH", ACCENT_RED),
        ("RULE-02", "Unknown Vehicle",
         "Vehicle not in KNOWN_VEHICLES list\nFires on new track entry only\nSeverity: MEDIUM", ACCENT_YLW),
        ("RULE-03", "Repeat Entry",
         "Same vehicle enters > REPEAT_ENTRY_LIMIT\nUses spatial track continuity\nSeverity: MEDIUM", ACCENT_ORG),
        ("RULE-04", "Person Loitering",
         "Activity = loitering or stationary\nFires on first sight + prolonged dwell\n(threshold: 15s test / 300s prod)\nSeverity: HIGH", ACCENT_RED),
        ("RULE-05", "Perimeter Activity",
         "Any object at 'perimeter' location\nduring after-hours window\nSeverity: HIGH", ACCENT_RED),
    ]

    for i, (rule_id, title, body, color) in enumerate(rules):
        x = 0.3 + i * 2.56
        add_component(slide, x, 1.1, 2.35, 2.8, f"{rule_id}\n{title}", body, color)

    # Flow note
    add_text_box(slide, "EVALUATION ORDER", 0.3, 4.2, 12.73, 0.3,
                 font_size=Pt(9), bold=True, color=GREY)
    flow = [
        ("Spatial Track\nUpdated", ACCENT_BLUE),
        ("RULE-01\nChecked", ACCENT_RED),
        ("RULE-02 & 03\nChecked", ACCENT_YLW),
        ("RULE-04\nChecked", ACCENT_RED),
        ("RULE-05\nChecked", ACCENT_RED),
        ("Alerts List\nReturned", ACCENT_GRN),
        ("Stored in\nSQLite", ACCENT_ORG),
        ("Telegram\nSent (HIGH)", ACCENT_GRN),
    ]
    fx = 0.3
    for label, color in flow:
        add_component(slide, fx, 4.55, 1.5, 0.9, "", label, color)
        fx += 1.5
        if fx < 0.3 + 1.5 * len(flow):
            add_arrow(slide, fx, 4.55 + 0.45, fx + 0.1, 4.55 + 0.45, color=GREY)
        fx += 0.12

    add_text_box(slide, "★  All HIGH alerts → Telegram photo notification  ·  MEDIUM alerts → SQLite only",
                 0.3, 6.1, 12.73, 0.35, font_size=Pt(9), color=GREY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 5 — Storage & Search
# ─────────────────────────────────────────────────────────────────────────────
def slide_storage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_text_box(slide, "DUAL-DATABASE STORAGE & SEMANTIC SEARCH", 0.3, 0.15, 12.73, 0.4,
                 font_size=Pt(14), bold=True, color=WHITE)

    # SQLite column
    add_text_box(slide, "SQLite  —  EventStore", 0.4, 0.75, 5.8, 0.35,
                 font_size=Pt(12), bold=True, color=ACCENT_RED)
    sqlite_items = [
        ("events table", "frame_id · timestamp · type · message · severity"),
        ("alerts table", "alert_id · rule_id · frame_id · message · severity"),
        ("Queries", "get_alerts(severity=) · get_events_by_timerange()"),
        ("Use case", "Structured lookups — time filters, alert history"),
    ]
    for i, (k, v) in enumerate(sqlite_items):
        y = 1.2 + i * 0.85
        add_component(slide, 0.4, y, 5.8, 0.75, k, v, ACCENT_RED)

    # ChromaDB column
    add_text_box(slide, "ChromaDB  —  FrameIndex", 7.0, 0.75, 5.8, 0.35,
                 font_size=Pt(12), bold=True, color=ACCENT_ORG)
    chroma_items = [
        ("Collection: frames", "Document = BLIP caption\nMetadata = frame_id, location, objects, bbox, track_id"),
        ("Embeddings", "sentence-transformers/all-MiniLM-L6-v2\n384-dim vectors · cosine similarity (HNSW)"),
        ("Queries", "query(text, n_results, location=, threat_level=)"),
        ("Use case", "Semantic search — find frames by meaning not keywords"),
    ]
    for i, (k, v) in enumerate(chroma_items):
        y = 1.2 + i * 0.85
        add_component(slide, 7.0, y, 5.8, 0.75, k, v, ACCENT_ORG)

    # Divider
    add_rect(slide, 6.5, 0.7, 0.04, 5.0, CARD_MID)

    # Query routing
    add_text_box(slide, "QUERY ENGINE — Routing Logic", 0.4, 4.72, 12.33, 0.35,
                 font_size=Pt(11), bold=True, color=ACCENT_BLUE)
    routing = [
        ('"how many people"', "→  ChromaDB list_all + track merge", ACCENT_BLUE),
        ('"after midnight"', "→  SQLite time-range filter", ACCENT_RED),
        ('"any alerts?"', "→  SQLite alerts table", ACCENT_RED),
        ('"person near door"', "→  ChromaDB semantic search", ACCENT_ORG),
    ]
    for i, (q, ans, color) in enumerate(routing):
        x = 0.4 + i * 3.2
        add_component(slide, x, 5.1, 3.0, 0.9, q, ans, color)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 6 — LangChain Agent
# ─────────────────────────────────────────────────────────────────────────────
def slide_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_text_box(slide, "LANGCHAIN REACT AGENT", 0.3, 0.15, 12.73, 0.4,
                 font_size=Pt(14), bold=True, color=WHITE)
    add_text_box(slide, "Context-aware event logging with GPT-4o-mini", 0.3, 0.52, 12.73, 0.3,
                 font_size=Pt(10), color=GREY)

    # Agent box
    add_component(slide, 4.5, 1.1, 4.33, 1.8,
                  "AgentExecutor",
                  "create_react_agent (LangChain 0.3)\nGPT-4o-mini via langchain-openai\nConversationBufferWindowMemory k=10\nmax_iterations=6  ·  handle_parsing_errors=True",
                  ACCENT_PRP, "🧠")

    # Tools
    tools = [
        ("log_event", "Records observation\nto SQLite events\ntable", ACCENT_BLUE),
        ("trigger_alert", "Creates alert entry\nin alerts table\n(agent-generated)", ACCENT_RED),
        ("query_history", "Semantic search of\nprevious frames\nvia ChromaDB", ACCENT_ORG),
    ]
    for i, (name, body, color) in enumerate(tools):
        x = 1.0 + i * 3.8
        add_component(slide, x, 3.4, 3.3, 1.6, f"@tool: {name}", body, color)
        add_arrow(slide, x + 1.65, 3.4, 6.67, 2.9, color=color, width=Pt(1.2))

    # Input
    add_component(slide, 0.3, 1.1, 3.8, 1.8,
                  "INPUT per frame",
                  "frame_id · timestamp · location\nobjects[ ] · activity\nBLIP caption · pre-fired alerts\nvehicle context (track_id, bbox)", ACCENT_BLUE, "📥")
    add_arrow(slide, 4.1, 2.0, 4.5, 2.0, color=ACCENT_BLUE, width=Pt(1.5))

    # Output
    add_component(slide, 9.1, 1.1, 3.93, 1.8,
                  "OUTPUT per frame",
                  "Natural language response\nstored as event log\nAgent memory carries context\nacross consecutive frames", ACCENT_GRN, "📤")
    add_arrow(slide, 8.83, 2.0, 9.1, 2.0, color=ACCENT_GRN, width=Pt(1.5))

    # ReAct loop note
    add_text_box(slide, "ReAct Loop: Thought → Action (tool call) → Observation → … → Final Answer",
                 0.3, 5.25, 12.73, 0.35, font_size=Pt(10), color=ACCENT_PRP, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "★  Deterministic rules fire FIRST — agent adds narrative context, not safety-critical decisions",
                 0.3, 5.65, 12.73, 0.35, font_size=Pt(9), color=GREY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 7 — Telegram Notification Flow
# ─────────────────────────────────────────────────────────────────────────────
def slide_telegram(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    add_text_box(slide, "TELEGRAM ALERT NOTIFICATION", 0.3, 0.15, 12.73, 0.4,
                 font_size=Pt(14), bold=True, color=WHITE)
    add_text_box(slide, "Real-time annotated frame delivery — HIGH severity only", 0.3, 0.52, 12.73, 0.3,
                 font_size=Pt(10), color=GREY)

    steps = [
        ("Alert Rule Fires\n(HIGH severity)", "RULE-01, RULE-04,\nor RULE-05 triggers\nin AlertRulesEngine", ACCENT_RED),
        ("telegram_alert()\nCalled", "rule_id · message\nseverity · frame_id\ntimestamp · location\nframe (numpy) · bbox", ACCENT_ORG),
        ("Frame Annotated\n(OpenCV)", "Copy of raw frame\nBounding box drawn\nRule label overlaid\nJPEG encoded (85%)", ACCENT_YLW),
        ("Daemon Thread\nLaunched", "Fire-and-forget\nNever blocks pipeline\n5s timeout on API", ACCENT_BLUE),
        ("sendPhoto\nAPI Call", "Bot API POST\nPhoto + caption\nPlain text format", ACCENT_GRN),
        ("Message Arrives\non Telegram", "🔴 Alert emoji\nSeverity · Frame #\nLocation · Message\n+ Annotated photo", ACCENT_GRN),
    ]

    for i, (title, body, color) in enumerate(steps):
        x = 0.35 + i * 2.16
        add_component(slide, x, 1.1, 1.95, 2.5, title, body, color)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.95, 2.35, x + 2.16, 2.35, color=GREY, width=Pt(1.5))

    # Disabled note
    add_component(slide, 0.35, 4.0, 5.8, 1.1,
                  "Graceful No-op",
                  "If TELEGRAM_BOT_TOKEN or TELEGRAM_CHAT_ID is empty\n→ send_alert() returns immediately, no error\nPipeline works fully without credentials", ACCENT_BLUE)

    add_component(slide, 6.8, 4.0, 5.8, 1.1,
                  "MEDIUM alerts suppressed",
                  "Only severity == 'high' reaches Telegram\nRULE-02 (unknown vehicle) and RULE-03 (repeat entry)\nare stored in SQLite but not notified", ACCENT_YLW)


# ─────────────────────────────────────────────────────────────────────────────
#  BUILD
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_pipeline(prs)
    slide_vlm(prs)
    slide_rules(prs)
    slide_storage(prs)
    slide_agent(prs)
    slide_telegram(prs)

    out = os.path.join(PROJECT_ROOT, "Flying_Police_Architecture.pptx")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
